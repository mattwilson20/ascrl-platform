from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import qrcode
import os

# ==============================
# CONFIGURATION
# ==============================
OUTPUT_FILE = "Veritas_International_Deck_v2.pptx"
CALENDLY_LINK = "https://calendly.com/dustin-blevins/15min"
LOGO_PATH = "logo.png"          # ← Replace with your logo
TRUCK_PHOTO = "truck.jpg"       # ← Replace with truck photo
DUSTIN_PHOTO = "dustin.jpg"     # ← Replace with Dustin headshot (optional)

# Colors
NAVY = RGBColor(10, 29, 86)      # #0A1D56
TEAL = RGBColor(0, 163, 163)     # #00A3A3
LIGHT_GRAY = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(220, 53, 69)
GREEN = RGBColor(40, 167, 69)

# Font sizes
TITLE_FS = 44
SUBTITLE_FS = 28
BODY_FS = 24
SMALL_FS = 18

# ==============================
# CREATE PRESENTATION
# ==============================
prs = Presentation()
prs.slide_width = Inches(13.33)  # 16:9
prs.slide_height = Inches(7.5)

# Set theme colors (approximate Montserrat)
def set_run_font(run, bold=False, size=Pt(24), color=NAVY):
    run.font.name = "Calibri"
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color

# Helper: Add title + subtitle
def add_title_subtitle(slide, title, subtitle=None):
    title_frame = slide.shapes.title
    title_frame.text = title
    set_run_font(title_frame.text_frame.paragraphs[0].runs[0], bold=True, size=Pt(TITLE_FS), color=WHITE)
    if subtitle:
        subtitle_tf = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(2))
        p = subtitle_tf.text_frame.paragraphs[0]
        p.text = subtitle
        p.alignment = PP_ALIGN.CENTER
        set_run_font(p.runs[0], bold=False, size=Pt(SUBTITLE_FS), color=WHITE)

# ==============================
# SLIDE 1: Cover
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
# Gradient background
fill = slide.background.fill
fill.gradient()
fill.gradient_stops[0].color.rgb = NAVY
fill.gradient_stops[1].color.rgb = TEAL

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(3))
tf = title_box.text_frame
tf.text = "Veritas International\nBuilt by Shippers for Shippers"
p = tf.paragraphs[0]
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER
tf.paragraphs[1].font.size = Pt(32)
tf.paragraphs[1].font.color.rgb = WHITE

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(10), Inches(1))
sub_box.text_frame.text = "Logistics Solutions from Dock to Door"
p = sub_box.text_frame.paragraphs[0]
p.font.size = Pt(28)
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Logo
if os.path.exists(LOGO_PATH):
    slide.shapes.add_picture(LOGO_PATH, Inches(10.5), Inches(0.3), height=Inches(1))

# Truck icon animation hint
note = slide.notes_slide.notes_text_frame
note.text = "Animation: Truck moves left → right on click"

# ==============================
# SLIDE 2: The Problem
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title + Content
slide.shapes.title.text = "TRUST IN 3PLs IS BROKEN"
set_run_font(slide.shapes.title.text_frame.paragraphs[0].runs[0], bold=True, size=Pt(TITLE_FS), color=NAVY)

# Table
left = Inches(0.5)
top = Inches(1.8)
width = Inches(12)
height = Inches(4)
table = slide.shapes.add_table(4, 2, left, top, width, height).table
table.cell(0, 0).text = "14-day average delay"
table.cell(0, 1).text = "+28% drayage costs in 2 years"
table.cell(1, 0).text = "No PO visibility"
table.cell(1, 1).text = "1 in 3 shipments hit demurrage"
table.cell(2, 0).text = "“Black box” pricing"
table.cell(2, 1).text = "Zero accountability"
table.cell(3, 0).text = "Source: 2025 Shippers Council Survey"
table.cell(3, 1).text = ""

for row in table.rows:
    for cell in row.cells:
        cell.text_frame.paragraphs[0].font.size = Pt(BODY_FS)
        cell.text_frame.paragraphs[0].font.color.rgb = NAVY if row.row_index < 3 else RGBColor(100,100,100)

# Icons
slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(1.9), Inches(0.4), Inches(0.4)).fill.solid().fore_color.rgb = RED
slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.9), Inches(0.4), Inches(0.4)).fill.solid().fore_color.rgb = RED

# ==============================
# SLIDE 3: Our Mission
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[5])
fill = slide.background.fill
fill.solid()
fill.fore_color.rgb = LIGHT_GRAY

words = ["TRANSPARENCY", "RESPECT", "HONESTY"]
top_start = Inches(1.5)
for i, word in enumerate(words):
    txbox = slide.shapes.add_textbox(Inches(1), top_start + i*Inches(1.2), Inches(11), Inches(1))
    p = txbox.text_frame.paragraphs[0]
    p.text = word
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = TEAL
    p.alignment = PP_ALIGN.CENTER

# Subtitle
sub = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11), Inches(1))
sub.text_frame.text = "We restore accountability.\nEvery shipment. Every time."
sub.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
set_run_font(sub.text_frame.paragraphs[0].runs[0], size=Pt(28), color=NAVY)

# Lock icon
lock = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.8), Inches(3), Inches(1.5), Inches(1.5))
lock.fill.fore_color.rgb = TEAL
lock.line.color.rgb = TEAL
lock.text_frame.text = "🔒"
lock.text_frame.paragraphs[0].font.size = Pt(60)
lock.text_frame.paragraphs[0].font.color.rgb = WHITE

# ==============================
# SLIDE 4: Global Reach
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "WCA-APPROVED NVOCC & FREIGHT FORWARDER"
content = slide.placeholders[1]
tf = content.text_frame
tf.clear()
items = [
    "11,000+ partners in 195 countries",
    "FMC-licensed, fully bonded",
    "FCL / LCL / Air consolidations",
    "Real-time PO visibility"
]
for item in items:
    p = tf.add_paragraph()
    p.text = "✓ " + item
    p.font.size = Pt(BODY_FS)
    p.font.color.rgb = NAVY

# Map placeholder
map_box = slide.shapes.add_textbox(Inches(7), Inches(2), Inches(5), Inches(4))
map_box.text_frame.text = "[WORLD MAP]\nInsert map with glowing dots:\nUSA • China • Europe"
map_box.text_frame.paragraphs[0].font.size = Pt(SMALL_FS)
map_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(100,100,100)

# ==============================
# SLIDE 5: Drayage Superpower
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "RABBIT TRANSPORT\nPowered by Veritas International"
set_run_font(slide.shapes.title.text_frame.paragraphs[0].runs[0], size=Pt(36), color=TEAL)

# Table
table = slide.shapes.add_table(4, 2, Inches(0.5), Inches(2), Inches(12), Inches(3)).table
left_col = ["Company-owned trucks & chassis", "<2-hour average gate turnaround", "98.7% on-time delivery (Q3)", "Proactive driver check-ins"]
right_col = ["Ports: Savannah | Atlanta | ARP", "Zero chassis splits in 2025", "", ""]
for i, (l, r) in enumerate(zip(left_col, right_col)):
    table.cell(i, 0).text = l
    table.cell(i, 1).text = r
    for cell in [table.cell(i,0), table.cell(i,1)]:
        cell.text_frame.paragraphs[0].font.size = Pt(BODY_FS - 2)
        cell.text_frame.paragraphs[0].font.color.rgb = NAVY

# Truck photo
if os.path.exists(TRUCK_PHOTO):
    slide.shapes.add_picture(TRUCK_PHOTO, Inches(8), Inches(3.5), width=Inches(4))

# ==============================
# SLIDE 6: Domestic Brokerage
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "NATIONWIDE CAPACITY. ZERO EXCUSES."
content = slide.placeholders[1].text_frame
content.clear()
services = [
    ("DRY VAN", "Full & partial loads", "AI-matched in <11 min"),
    ("OPEN DECK", "Flatbed | Step-deck | RGN", "Heavy & OD up to 120K lbs"),
    ("EXPEDITED", "24/7 hotshot & team", "Door-to-door in 18 hrs")
]
for icon, title, desc in services:
    p = content.add_paragraph()
    p.text = f"{title}\n{desc}"
    p.font.size = Pt(BODY_FS)
    p.font.bold = True if "VAN" in title or "DECK" in title or "EXPEDITED" in title else False

# Bottom bullets
p = content.add_paragraph()
p.text = "✓ 98% carrier retention\n✓ No fuel surcharge markups – ever"
p.font.size = Pt(BODY_FS)
p.font.color.rgb = GREEN

# ==============================
# SLIDE 7: Proof
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "PROOF IN EVERY MILE"
quote = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(10), Inches(1.5))
quote.text_frame.text = "“Veritas saved us $47K in demurrage last quarter.”\n– John D., ACME Imports"
quote.text_frame.paragraphs[0].font.size = Pt(26)
quote.text_frame.paragraphs[0].font.italic = True
quote.text_frame.paragraphs[0].font.color.rgb = NAVY

# Stats circles
stats = [("98.7%", "ON-TIME"), ("<2 hr", "GATE TURN"), ("Zero", "CHASSIS SPLITS")]
left = Inches(2)
for i, (num, label) in enumerate(stats):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + i*Inches(3), Inches(4), Inches(1.8), Inches(1.8))
    circle.fill.fore_color.rgb = TEAL
    circle.line.color.rgb = TEAL
    tf = circle.text_frame
    tf.text = num + "\n" + label
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[1].font.size = Pt(16)

# ==============================
# SLIDE 8: Let’s Move
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[5])
fill = slide.background.fill
fill.solid()
fill.fore_color.rgb = WHITE

if os.path.exists(DUSTIN_PHOTO):
    slide.shapes.add_picture(DUSTIN_PHOTO, Inches(0.5), Inches(0.5), height=Inches(6))
    pic = slide.shapes[-1]
    pic._element.get_or_add_ln().add_prstDash(val='solid')

# Contact card
card = slide.shapes.add_textbox(Inches(7), Inches(1), Inches(5.5), Inches(5))
tf = card.text_frame
tf.text = "LET’S MOVE TOGETHER\n\nDustin Blevins\nCo-Founder & President\n\n📞 423-480-0085\n✉️ Dblevins@jcfvi.com\n🌐 jcfvi.com"
for p in tf.paragraphs:
    p.font.size = Pt(24) if p.text.startswith("LET") else Pt(20)
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.LEFT

# QR Code
qr = qrcode.QRCode(version=1, box_size=10, border=2)
qr.add_data(CALENDLY_LINK)
qr.make(fit=True)
img = qr.make_image(fill_color="navy", back_color="white")
qr_path = "calendly_qr.png"
img.save(qr_path)
slide.shapes.add_picture(qr_path, Inches(9.5), Inches(5), height=Inches(1.5))

# Quotes
quote1 = slide.shapes.add_textbox(Inches(7), Inches(6.2), Inches(5.5), Inches(1))
quote1.text_frame.text = "“A company who rewards loyalty.”\n“A company who makes lives better.”"
quote1.text_frame.paragraphs[0].font.size = Pt(18)
quote1.text_frame.paragraphs[0].font.italic = True

# Cleanup QR
os.remove(qr_path)

# ==============================
# SLIDE 9: Thank You (Backup)
# ==============================
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5)).text_frame.text = \
    "THANK YOU\n\n[Backup Stats]\n• 2025 YTD: 2,847 containers moved\n• Avg. savings per TEU: $183\n• Customer NPS: 89\n\nNeed rates? Text “BOL” to 423-480-0085"

# Hide backup slide
slide_part = slide.part
slide_part._element.set('show', '0')  # Hide in show

# ==============================
# SAVE
# ==============================
prs.save(OUTPUT_FILE)
print(f"✅ Deck created: {OUTPUT_FILE}")
print("   Open in PowerPoint → Add real logo, truck, Dustin photo → Present!")